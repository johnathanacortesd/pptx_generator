# ============================================
# KREAB COLOMBIA - SUITE DE INTELIGENCIA ELECTORAL
# Streamlit App v1.2 - Cache + Auth + Icons
# ============================================

import streamlit as st

st.set_page_config(
    page_title="Kreab - Inteligencia Electoral",
    page_icon="🗳️📊",
    layout="wide",
    initial_sidebar_state="expanded"
)

import pandas as pd
import numpy as np
import networkx as nx
import plotly.graph_objects as go
from collections import Counter, OrderedDict
import re
import math
import random
from datetime import datetime
from wordcloud import WordCloud
import matplotlib.pyplot as plt
import matplotlib
matplotlib.use('Agg')
import io
import base64
import json

try:
    from unidecode import unidecode
except:
    unidecode = lambda x: x

try:
    from openai import OpenAI
    OPENAI_AVAILABLE = True
except:
    OPENAI_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    PPTX_AVAILABLE = True
except:
    PPTX_AVAILABLE = False

KALEIDO_OK = False
try:
    test_fig = go.Figure(data=[go.Scatter(x=[1,2], y=[1,2])])
    test_bytes = test_fig.to_image(format="png", width=100, height=100)
    if test_bytes and len(test_bytes) > 100:
        KALEIDO_OK = True
except:
    pass

# ============================================
# AUTENTICACIÓN
# ============================================

def check_password():
    if "authenticated" not in st.session_state:
        st.session_state.authenticated = False
    if st.session_state.authenticated:
        return True

    st.markdown("""
    <div style="background:linear-gradient(135deg,#8B0000,#111827);color:white;padding:28px;border-radius:14px;margin-bottom:22px;text-align:center">
        <h1 style="margin:0">🗳️📊 PPTX_JC Intelligence</h1>
        <p style="margin:5px 0 0;opacity:0.9">Acceso restringido</p>
    </div>
    """, unsafe_allow_html=True)

    col1, col2, col3 = st.columns([1, 2, 1])
    with col2:
        with st.form("login_form"):
            password = st.text_input("🔑 Contraseña", type="password", placeholder="Ingresa la contraseña")
            login_btn = st.form_submit_button("Ingresar", use_container_width=True, type="primary")

            if login_btn:
                try:
                    correct_password = st.secrets["APP_PASSWORD"]
                except:
                    st.error("❌ APP_PASSWORD no configurada en Secrets")
                    return False
                if password == correct_password:
                    st.session_state.authenticated = True
                    st.rerun()
                else:
                    st.error("❌ Contraseña incorrecta")
                    return False
    return False

# ============================================
# CONFIGURACIÓN
# ============================================

COLORES = {
    'Mauricio Cárdenas': '#B22222',
    'Juan Carlos Pinzón': '#00BFFF',
    'Paloma Valencia': '#00008B',
    'Juan Daniel Oviedo': '#D2691E',
    'Juan Manuel Galán': '#9400D3'
}

CANDIDATOS_CONFIG = {
    'Mauricio Cárdenas': {
        'nombre_corto': 'M. Cárdenas',
        'regex': [r'mauricio\s+c[aá]rdenas', r'\bc[aá]rdenas\b'],
        'blacklist': ['mauricio','cardenas','cárdenas','ministro','doctor','mcardenas','santos']
    },
    'Juan Carlos Pinzón': {
        'nombre_corto': 'JC Pinzón',
        'regex': [r'juan\s+carlos\s+pinz[oó]n', r'\bpinz[oó]n\b', r'jc\s*pinzon'],
        'blacklist': ['juan','carlos','pinzon','pinzón','jcpinzon','ministro','embajador','jc']
    },
    'Paloma Valencia': {
        'nombre_corto': 'P. Valencia',
        'regex': [r'paloma\s+valencia', r'\bpaloma\b', r'\bvalencia\b'],
        'blacklist': ['paloma','valencia','senadora','pvalencia','democrata','demócrata','centro']
    },
    'Juan Daniel Oviedo': {
        'nombre_corto': 'JD Oviedo',
        'regex': [r'juan\s+daniel\s+oviedo', r'\boviedo\b', r'jd\s*oviedo'],
        'blacklist': ['juan','daniel','oviedo','director','dane','jdoviedo','jd']
    },
    'Juan Manuel Galán': {
        'nombre_corto': 'JM Galán',
        'regex': [r'juan\s+manuel\s+gal[aá]n', r'\bgal[aá]n\b'],
        'blacklist': ['juan','manuel','galan','galán','jmgalan','jm','nuevo','liberalismo','luis','carlos']
    }
}

def construir_blacklist_global_candidatos():
    bl = set()
    for cand, cfg in CANDIDATOS_CONFIG.items():
        for w in cfg['blacklist']:
            bl.add(w.lower().strip());bl.add(unidecode(w.lower().strip()))
        for p in re.findall(r'[A-Za-zÁÉÍÓÚÑáéíóúñ]+', cand.lower()):
            bl.add(p);bl.add(unidecode(p))
        for p in re.findall(r'[A-Za-zÁÉÍÓÚÑáéíóúñ]+', cfg['nombre_corto'].lower()):
            bl.add(p);bl.add(unidecode(p))
        iniciales = "".join([p[0] for p in re.findall(r'[a-záéíóúñ]+', cand.lower()) if p])
        if len(iniciales) >= 2: bl.add(iniciales)
    return bl

BLACKLIST_GLOBAL_CANDIDATOS = construir_blacklist_global_candidatos()

# ============================================
# CLIENTE OPENAI
# ============================================

def obtener_cliente_openai():
    if not OPENAI_AVAILABLE: return None
    try:
        api_key = st.secrets.get("OPENAI_API_KEY", None)
        if not api_key: return None
        return OpenAI(api_key=api_key)
    except:
        return None

# ============================================
# NORMALIZACIÓN Y PROCESAMIENTO
# ============================================

def normalizar_palabra(palabra):
    if pd.isna(palabra): return ""
    palabra = unidecode(str(palabra).lower().strip())
    if palabra.endswith('es') and len(palabra) > 4: palabra = palabra[:-2]
    elif palabra.endswith('s') and len(palabra) > 4 and not palabra.endswith('ss'): palabra = palabra[:-1]
    return palabra

def normalizar_token(token):
    if not token: return ""
    token = str(token).strip().lower()
    if token.startswith("#"):
        core = unidecode(token[1:]);core = re.sub(r'[^a-z0-9_ñ]+','',core)
        return "#"+core if len(core)>=3 else ""
    token = re.sub(r'[^a-záéíóúñ]+','',token.lower())
    return normalizar_palabra(token) if len(token)>=4 else ""

def normalizar_token_nube(token):
    if not token: return ""
    token = str(token).strip().lower()
    if token.startswith("#"):
        core = unidecode(token[1:]);core = re.sub(r'[^a-z0-9_]+','',core)
        return "#"+core if len(core)>=3 else ""
    token_norm = unidecode(token);token_norm = re.sub(r'[^a-z]+','',token_norm)
    if len(token_norm) < 4: return ""
    if token_norm.endswith('es') and len(token_norm) > 5: token_norm = token_norm[:-2]
    elif token_norm.endswith('s') and len(token_norm) > 5 and not token_norm.endswith('ss'): token_norm = token_norm[:-1]
    return token_norm

def limpiar_texto(texto):
    if pd.isna(texto): return ""
    return str(texto).lower()

def identificar_candidatos(texto):
    texto = limpiar_texto(texto)
    return list(set([cand for cand,cfg in CANDIDATOS_CONFIG.items()
                     if any(re.search(pat,texto) for pat in cfg['regex'])]))

STOPWORDS_BASE = set([
    'el','la','los','las','un','una','de','del','y','o','en','a','para','por','con','su','sus','anos',
    'es','son','fue','ha','al','lo','se','que','qué','como','más','pero','no','si','ya','todo','nada',
    'esta','este','rt','via','amp','https','http','tco','t','co','com','www','status','twitter','x','quot',
    'sobre','cuando','porque','colombia','decide','candidato','presidente','donde','hace','entre','detall',
    'gran','muy','desde','todos','quien','pais','país','gente','política','solo','ahora','vamos','jajaja','jajajaj',
    'tiene','tienen','esto','eso','estos','esas','uno','dos','tres','cada','ser','puede','sido','tinyur',
    'debe','están','hacer','tras','otras','bien','aquí','hasta','hoy','mañana','ayer','asi','aun',
    'tambien','también','pues','sino','toda','todas','todo','todos','tanto','tan','porque'
])

def obtener_tokens_filtrados(series_textos, blacklist_extra=None, incluir_hashtags=True):
    blacklist_extra = set([normalizar_token(x) for x in (blacklist_extra or []) if x])
    full_text = " ".join([str(t) for t in series_textos if not pd.isna(t)]).lower()
    tokens = []
    if incluir_hashtags: tokens.extend(re.findall(r'#\w+', full_text))
    tokens.extend(re.findall(r'\b[a-záéíóúñ]{4,}\b', full_text))
    filtered = []
    for t in [normalizar_token(t) for t in tokens]:
        if not t: continue
        if t.startswith("#"):
            if t not in blacklist_extra and t not in {'#colombiadecide2026','#debate','#colombia'}: filtered.append(t)
        else:
            if t not in STOPWORDS_BASE and t not in blacklist_extra: filtered.append(t)
    return filtered

def obtener_tokens_nube(series_textos, blacklist_extra=None):
    blacklist = set()
    if blacklist_extra:
        for w in blacklist_extra:
            if w:
                bl_norm = unidecode(str(w).lower().strip())
                blacklist.add(bl_norm);blacklist.add(normalizar_token_nube(w))
    for nombre_part in BLACKLIST_GLOBAL_CANDIDATOS:
        blacklist.add(nombre_part);blacklist.add(unidecode(nombre_part));blacklist.add(normalizar_token_nube(nombre_part))
    full_text = " ".join([str(t) for t in series_textos if not pd.isna(t)]).lower()
    hashtags_raw = re.findall(r'#\w+', full_text);words_raw = re.findall(r'\b[a-záéíóúñ]{4,}\b', full_text)
    freq_map = {}
    for tag in hashtags_raw:
        norm = normalizar_token_nube(tag)
        if not norm: continue
        if norm.replace('#','') in blacklist: continue
        if norm in {'#colombiadecide2026','#debate','#colombia'}: continue
        freq_map[norm] = freq_map.get(norm, 0) + 1
    stopwords_norm = set([unidecode(s) for s in STOPWORDS_BASE])
    for word in words_raw:
        norm = normalizar_token_nube(word)
        if not norm or norm in stopwords_norm or norm in blacklist: continue
        freq_map[norm] = freq_map.get(norm, 0) + 1
    return freq_map

def _short_label(s, maxlen=18):
    s = str(s);return s if len(s)<=maxlen else s[:maxlen-1]+"…"

# ============================================
# ORDENAR
# ============================================

def ordenar_kpis_por_menciones(kpis):
    return OrderedDict(sorted(kpis.items(), key=lambda x: x[1]['menciones'], reverse=True))

def ordenar_ejes_por_menciones_interacciones(ejes_tematicos, kpis):
    orden = sorted(kpis.keys(), key=lambda c: kpis[c]['menciones'], reverse=True)
    result = OrderedDict()
    for cand in orden:
        if cand not in ejes_tematicos or not ejes_tematicos[cand]: continue
        eje = ejes_tematicos[cand].copy();temas = []
        for i in range(1,6):
            tk = f'tema_{i}'
            if tk in eje and eje[tk]:
                tema = eje[tk].copy();tema['_score'] = tema.get('posts',0)+tema.get('interacciones',0);temas.append(tema)
        temas.sort(key=lambda t: t['_score'], reverse=True)
        en = {'total_posts':eje.get('total_posts',0),'total_interacciones':eje.get('total_interacciones',0)}
        if 'fallback' in eje: en['fallback']=eje['fallback']
        for i,t in enumerate(temas,1): t.pop('_score',None); en[f'tema_{i}']=t
        result[cand] = en
    return result

def normalizar_temas_entre_candidatos(ejes_tematicos, client=None):
    todos_temas = []
    for cand,eje in ejes_tematicos.items():
        if not eje: continue
        for i in range(1,6):
            tk=f'tema_{i}'
            if tk in eje and eje[tk]:
                nombre=eje[tk].get('nombre','')
                if nombre: todos_temas.append({'candidato':cand,'tema_key':tk,'nombre_original':nombre})
    if not todos_temas: return ejes_tematicos
    nombres_unicos = list(set([t['nombre_original'] for t in todos_temas]))
    if len(nombres_unicos) <= 1: return ejes_tematicos
    if client:
        try:
            lista = "\n".join([f"- {n}" for n in nombres_unicos])
            prompt = f"""Lista de temas electorales colombianos:
{lista}
Normaliza SOLO sinonimos exactos. JSON: {{"nombre original": "nombre normalizado"}}"""
            response = client.chat.completions.create(model="gpt-4.1-nano-2025-04-14",
                messages=[{"role":"system","content":"Solo sinonimos exactos. JSON valido."},{"role":"user","content":prompt}],
                temperature=0.05, max_tokens=600)
            content = response.choices[0].message.content.strip()
            if content.startswith("```"):
                content = re.sub(r'^```json?\s*','',content);content = re.sub(r'\s*```$','',content)
            mapeo = json.loads(content)
            for ti in todos_temas:
                ejes_tematicos[ti['candidato']][ti['tema_key']]['nombre'] = mapeo.get(ti['nombre_original'],ti['nombre_original'])
            return ejes_tematicos
        except: pass
    norm_map = {}
    for ti in todos_temas:
        key = unidecode(ti['nombre_original'].lower().strip())
        if key not in norm_map: norm_map[key] = ti['nombre_original']
    for ti in todos_temas:
        key = unidecode(ti['nombre_original'].lower().strip())
        ejes_tematicos[ti['candidato']][ti['tema_key']]['nombre'] = norm_map[key]
    return ejes_tematicos

def sanitizar_nombre_para_prompt(nombre): return unidecode(str(nombre))

def parsear_json_ia(content):
    content = content.strip()
    if content.startswith("```"):
        content = re.sub(r'^```json?\s*','',content);content = re.sub(r'\s*```$','',content)
    content = content.strip()
    try: return json.loads(content)
    except: pass
    try: return json.loads(content.replace("'",'"'))
    except: pass
    try:
        m = re.search(r'\{[\s\S]*\}', content)
        if m: return json.loads(m.group(0))
    except: pass
    try: return json.loads(unidecode(content))
    except: pass
    raise json.JSONDecodeError("No parse", content, 0)

# ============================================
# EJES TEMÁTICOS
# ============================================

def analizar_ejes_tematicos_con_ia(df_exploded, candidato, client):
    subset = df_exploded[df_exploded['Candidatos_Detectados']==candidato].copy()
    if subset.empty: return None
    blacklist = CANDIDATOS_CONFIG[candidato]['blacklist']
    tokens = obtener_tokens_filtrados(subset['Contenido de la publicación'], blacklist, incluir_hashtags=True)
    freq = Counter(tokens).most_common(40);ci = None
    for col in ['Interacciones totales','interacciones totales','Interacciones','interactions']:
        if col in subset.columns: ci=col; break
    ti = int(subset[ci].fillna(0).sum()) if ci else len(subset);tp = len(subset)
    terms = []
    for term,count in freq[:30]:
        mask = subset['Contenido de la publicación'].str.lower().str.contains(term,na=False,regex=False)
        inter = int(subset.loc[mask,ci].fillna(0).sum()) if ci and mask.any() else count
        terms.append(f"{term}: {count} posts, {inter:,} interacciones")
    sample = "\n".join([f"- {str(c)[:180]}" for c in subset['Contenido de la publicación'].dropna().head(20)])
    candidato_safe = sanitizar_nombre_para_prompt(candidato)
    prompt = f"""Candidato {candidato_safe}, campana Colombia 2026.
FRECUENCIA: {chr(10).join(terms[:25])}
MUESTRA ({tp} posts): {sample}
METRICAS: {tp} posts, {ti:,} interacciones
5 TEMAS ESPECIFICOS. ASCII only. Solo JSON valido:
{{"tema_1":{{"nombre":"Tema","posts":n,"interacciones":n,"palabras_clave":["p1","p2","p3"]}},"tema_2":{{"nombre":"Tema","posts":n,"interacciones":n,"palabras_clave":["p1","p2","p3"]}},"tema_3":{{"nombre":"Tema","posts":n,"interacciones":n,"palabras_clave":["p1","p2","p3"]}},"tema_4":{{"nombre":"Tema","posts":n,"interacciones":n,"palabras_clave":["p1","p2","p3"]}},"tema_5":{{"nombre":"Tema","posts":n,"interacciones":n,"palabras_clave":["p1","p2","p3"]}}}}"""
    try:
        r = client.chat.completions.create(model="gpt-4.1-nano-2025-04-14",
            messages=[{"role":"system","content":"Analista politico. JSON valido. Sin acentos."},{"role":"user","content":prompt}],
            temperature=0.3, max_tokens=800)
        resultado = parsear_json_ia(r.choices[0].message.content.strip())
        resultado['total_posts'] = tp; resultado['total_interacciones'] = ti;return resultado
    except:
        return generar_analisis_fallback(subset,freq,candidato,tp,ti,ci)

def generar_analisis_fallback(subset,freq,candidato,total_posts,total_inter,col_inter):
    tk = {'Seguridad':['seguridad','policia','crimen','violencia','paz'],'Economía':['economia','empleo','trabajo','impuestos','reforma'],
          'Corrupción':['corrupcion','robo','investigacion','fiscalia'],'Educación':['educacion','universidad','escuela'],
          'Salud':['salud','hospital','eps'],'Elecciones':['elecciones','votos','campana','debate','encuesta'],
          'Gobierno':['gobierno','petro','oposicion','congreso'],'Justicia':['justicia','juez','tribunal'],
          'Medio ambiente':['ambiente','clima','mineria'],'Infraestructura':['infraestructura','vias','transporte']}
    tc={};tid={}
    for tema,kws in tk.items():
        c=0;it=0
        for term,fc in freq:
            if any(kw in term.replace('#','').lower() for kw in kws):
                c+=fc
                if total_posts>0:it+=int((fc/total_posts)*total_inter)
        if c>0:tc[tema]=c;tid[tema]=it
    to=sorted(tc.keys(),key=lambda x:tc[x]+tid.get(x,0),reverse=True)[:5]
    if len(to)<5:to.extend(['Política general','Debate electoral','Propuestas','Candidatura','Campaña'][:5-len(to)])
    r={'total_posts':total_posts,'total_interacciones':total_inter,'fallback':True}
    for i,t in enumerate(to[:5],1):
        r[f'tema_{i}']={'nombre':t,'posts':tc.get(t,int(total_posts/5)),'interacciones':tid.get(t,int(total_inter/5)),'palabras_clave':[x[0] for x in freq[:3]]}
    return r

# ============================================
# LAYOUT GEPHI
# ============================================

def _hex_to_rgb(hc):
    hc=hc.lstrip('#');return tuple(int(hc[i:i+2],16) for i in (0,2,4))

def _mezclar_con_blanco(rgb,fb):
    r,g,b=rgb;return(int(r*(1-fb)+255*fb),int(g*(1-fb)+255*fb),int(b*(1-fb)+255*fb))

def construir_radios_nodos(G):
    radii={}
    for n,d in G.nodes(data=True):
        t=d.get('type','palabra');l=str(d.get('label',n));w=float(d.get('weight',1) or 1)
        if t=='candidato':r=7.0
        elif t=='autor':r=2.0+0.15*min(len(l),22)+0.12*math.log1p(w)
        else:r=1.8+0.12*min(len(l),18)+0.10*math.log1p(w)
        radii[n]=r
    return radii

def resolver_colisiones(pos,radii,iters=150,padding=0.8):
    nodes=list(pos.keys())
    for _ in range(iters):
        moved=0
        for i in range(len(nodes)):
            for j in range(i+1,len(nodes)):
                n1,n2=nodes[i],nodes[j];p1,p2=pos[n1],pos[n2]
                dx,dy=p2[0]-p1[0],p2[1]-p1[1];dist=math.hypot(dx,dy)
                min_d=radii.get(n1,1.5)+radii.get(n2,1.5)+padding
                if dist<1e-9:ang=random.random()*2*math.pi;dx,dy=math.cos(ang),math.sin(ang);dist=1e-6
                if dist<min_d:
                    ux,uy=dx/dist,dy/dist;s=(min_d-dist)/2.0*1.1
                    pos[n1]=(p1[0]-ux*s,p1[1]-uy*s);pos[n2]=(p2[0]+ux*s,p2[1]+uy*s);moved+=1
        if moved==0:break
    return pos

def layout_force_atlas_mejorado(G,iteraciones=600,gravity=0.012,scaling_ratio=22.0,seed=42,radii=None):
    random.seed(seed);np.random.seed(seed)
    radii=radii or construir_radios_nodos(G);pos={}
    cands=[n for n,d in G.nodes(data=True) if d.get('type')=='candidato']
    nc=max(1,len(cands));R=20.0
    for i,c in enumerate(cands):ang=(2*math.pi*i)/nc;pos[c]=np.array([R*math.cos(ang),R*math.sin(ang)],dtype=float)
    for n,d in G.nodes(data=True):
        if n in pos:continue
        p=d.get('parent')
        if p and p in pos:pos[n]=pos[p]+np.array([random.uniform(-8,8),random.uniform(-8,8)],dtype=float)
        else:pos[n]=np.array([random.uniform(-12,12),random.uniform(-12,12)],dtype=float)
    nodes=list(G.nodes());deg=dict(G.degree())
    for it in range(iteraciones):
        disp={n:np.array([0.0,0.0],dtype=float) for n in nodes}
        for i in range(len(nodes)):
            for j in range(i+1,len(nodes)):
                n1,n2=nodes[i],nodes[j];delta=pos[n2]-pos[n1];dist=float(np.linalg.norm(delta))+1e-9
                md=radii.get(n1,1.5)+radii.get(n2,1.5)+0.5
                if dist<md:dist=md
                f=(scaling_ratio*(deg.get(n1,1)+1)*(deg.get(n2,1)+1))/(dist**2)
                d2=delta/dist;disp[n1]-=d2*f;disp[n2]+=d2*f
        for u,v,data in G.edges(data=True):
            delta=pos[v]-pos[u];dist=float(np.linalg.norm(delta))+1e-9;w=float(data.get('weight',1.0) or 1.0)
            f=(dist**2)/(scaling_ratio)*(0.05+0.015*math.log1p(w));d2=delta/dist;disp[u]+=d2*f;disp[v]-=d2*f
        for n in nodes:p=pos[n];dc=float(np.linalg.norm(p))+1e-9;disp[n]-=(p/dc)*(gravity*dc)
        sp=0.10+0.90*(1-it/iteraciones);ms=2.2*sp
        for n in nodes:
            step=disp[n]*sp;mag=float(np.linalg.norm(step))
            if mag>ms:step=step/mag*ms
            pos[n]+=step
        if it%6==0:
            pt={n:(float(pos[n][0]),float(pos[n][1])) for n in nodes}
            pt=resolver_colisiones(pt,radii,iters=12,padding=0.9)
            for n in nodes:pos[n]=np.array([pt[n][0],pt[n][1]],dtype=float)
    pf={n:(float(pos[n][0]),float(pos[n][1])) for n in nodes}
    return resolver_colisiones(pf,radii,iters=200,padding=1.0)

# ============================================
# REDES
# ============================================

def crear_red_principal(df_exploded):
    G=nx.Graph()
    for cand in CANDIDATOS_CONFIG:
        G.add_node(cand,type='candidato',label=CANDIDATOS_CONFIG[cand]['nombre_corto'],color=COLORES[cand],weight=1000)
        sub=df_exploded[df_exploded['Candidatos_Detectados']==cand]
        if sub.empty:continue
        bl=CANDIDATOS_CONFIG[cand]['blacklist'];toks=obtener_tokens_filtrados(sub['Contenido de la publicación'],bl,incluir_hashtags=True)
        for tok,freq in Counter(toks).most_common(15):
            nid=f"t_{cand}_{tok}";G.add_node(nid,type='palabra',label=tok,parent=cand,color=COLORES[cand],weight=int(freq));G.add_edge(cand,nid,weight=int(freq))
        if 'Autor' in sub.columns:
            for aut,freq in sub['Autor'].value_counts().head(6).items():
                if pd.isna(aut):continue
                nid=f"a_{cand}_{aut}";G.add_node(nid,type='autor',label=f"@{_short_label(aut,14)}",label_full=f"@{aut}",parent=cand,color=COLORES[cand],weight=int(freq));G.add_edge(cand,nid,weight=int(freq))
    radii=construir_radios_nodos(G);pos=layout_force_atlas_mejorado(G,iteraciones=650,gravity=0.028,scaling_ratio=19.0,seed=7,radii=radii)
    fig=go.Figure()
    ex,ey=[],[]
    for u,v,_ in G.edges(data=True):
        if u in pos and v in pos:ex+=[pos[u][0],pos[v][0],None];ey+=[pos[u][1],pos[v][1],None]
    fig.add_trace(go.Scatter(x=ex,y=ey,mode='lines',line=dict(color='rgba(180,180,180,0.45)',width=1.0),hoverinfo='none',showlegend=False))
    cx,cy,ct,cc=[],[],[],[]
    for n,d in G.nodes(data=True):
        if d.get('type')=='candidato':cx.append(pos[n][0]);cy.append(pos[n][1]);ct.append(d.get('label',n));cc.append(d.get('color','#333'))
    fig.add_trace(go.Scatter(x=cx,y=cy,mode='markers+text',marker=dict(size=58,color=cc,line=dict(width=4,color='white')),text=ct,textposition='bottom center',textfont=dict(color='black',size=18,family='Segoe UI, Arial Black'),hoverinfo='none',showlegend=False))
    wx,wy,ws,wc,wh,wsym=[],[],[],[],[],[]
    for n,d in G.nodes(data=True):
        t=d.get('type')
        if t=='candidato':continue
        x,y=pos[n];w=int(d.get('weight',1) or 1);wx.append(x);wy.append(y)
        if t=='autor':wsym.append('diamond');ws.append(min(10+2*math.log1p(w),18));wc.append('rgba(90,90,90,0.9)');wh.append(f"<b>{d.get('label_full',d.get('label'))}</b><br>Autor<br>{w}")
        else:wsym.append('circle');ws.append(min(12+2.5*math.log1p(w),24));wc.append('rgba(70,70,70,0.9)');wh.append(f"<b>{d.get('label')}</b><br>Palabra/Hashtag<br>{w}")
    fig.add_trace(go.Scatter(x=wx,y=wy,mode='markers',marker=dict(size=ws,color=wc,symbol=wsym,line=dict(width=1,color='white')),hovertemplate="%{hovertext}<extra></extra>",hovertext=wh,showlegend=False))
    annotations=[]
    for n,d in G.nodes(data=True):
        if d.get('type')=='candidato' or n not in pos:continue
        t=d.get('type');label=str(d.get('label',n));fs=12 if t=='palabra' else 11
        if len(label)>=14:fs=max(7,fs-1)
        border=d.get('color','#666') if t in ('palabra','autor') else '#666'
        annotations.append(dict(x=pos[n][0],y=pos[n][1],text=f"<b>{label}</b>",font=dict(color='white',size=fs,family='Segoe UI, Arial'),bgcolor="rgba(0,0,0,0.85)",bordercolor=border,borderwidth=1,borderpad=2,opacity=1,showarrow=False))
    xs=[p[0] for p in pos.values()];ys=[p[1] for p in pos.values()];pad=10
    fig.update_layout(annotations=annotations,showlegend=False,plot_bgcolor='white',margin=dict(t=10,b=10,l=10,r=10),xaxis=dict(showgrid=False,zeroline=False,visible=False,range=[min(xs)-pad,max(xs)+pad]),yaxis=dict(showgrid=False,zeroline=False,visible=False,range=[min(ys)-pad,max(ys)+pad]),height=1200)
    return fig,G,pos

def crear_red_inversa_gephi(df_exploded):
    G=nx.Graph();all_bl=[]
    for cfg in CANDIDATOS_CONFIG.values():all_bl.extend(cfg['blacklist'])
    tg=obtener_tokens_filtrados(df_exploded['Contenido de la publicación'],all_bl,incluir_hashtags=True);top_g=[w for w,c in Counter(tg).most_common(35)];pc={}
    for word in top_g:
        ct=0;cm=[]
        for cand in CANDIDATOS_CONFIG:
            sub=df_exploded[df_exploded['Candidatos_Detectados']==cand];txt=" ".join(sub['Contenido de la publicación'].astype(str)).lower()
            c=txt.count(word) if word.startswith("#") else normalizar_palabra(txt).count(word)
            if c>0:ct+=c;cm.append((cand,c))
            if c>0 and not G.has_node(cand):G.add_node(cand,type='candidato',label=CANDIDATOS_CONFIG[cand]['nombre_corto'],color=COLORES[cand])
            if c>0:G.add_edge(cand,word,weight=c)
        if ct>0:G.add_node(word,type='palabra',label=word,weight=ct);pc[word]=cm
    radii=construir_radios_nodos(G);pos=layout_force_atlas_mejorado(G,iteraciones=350,gravity=0.025,scaling_ratio=18.0,seed=11,radii=radii)
    fig=go.Figure()
    for u,v,data in G.edges(data=True):
        w=data.get('weight',1);width=min(1+(w/25),3.5)
        fig.add_trace(go.Scatter(x=[pos[u][0],pos[v][0]],y=[pos[u][1],pos[v][1]],mode='lines',line=dict(color='rgba(200,200,200,0.65)',width=width),hoverinfo='none',showlegend=False))
    cx,cy,ct,cc,cs=[],[],[],[],[]
    for n,d in G.nodes(data=True):
        if d['type']=='candidato':cx.append(pos[n][0]);cy.append(pos[n][1]);ct.append(d['label']);cc.append(d['color']);cs.append(55)
    fig.add_trace(go.Scatter(x=cx,y=cy,mode='markers+text',marker=dict(size=cs,color=cc,line=dict(width=4,color='white')),text=ct,textposition='bottom center',textfont=dict(color='black',size=17,family='Segoe UI, Arial Black'),hoverinfo='none',showlegend=False))
    wx,wy,wh,ws=[],[],[],[]
    for n,d in G.nodes(data=True):
        if d['type']=='palabra':
            wx.append(pos[n][0]);wy.append(pos[n][1]);ws.append(min(10+(d['weight']/6),26))
            ci=pc.get(n,[]);ct2="<br>".join([f"  • {CANDIDATOS_CONFIG[c]['nombre_corto']}: {f}" for c,f in sorted(ci,key=lambda x:x[1],reverse=True)])
            wh.append(f"<b>{d['label']}</b><br>Menciones: {d['weight']}<br><br>{ct2}")
    fig.add_trace(go.Scatter(x=wx,y=wy,mode='markers',marker=dict(size=ws,color='#555',line=dict(width=1,color='white')),hovertemplate="%{hovertext}<extra></extra>",hovertext=wh,showlegend=False))
    annotations=[]
    for n,d in G.nodes(data=True):
        if d['type']=='palabra':annotations.append(dict(x=pos[n][0],y=pos[n][1],text=f"<b>{d['label']}</b>",font=dict(color='white',size=15,family='Segoe UI, Arial'),bgcolor="rgba(0,0,0,0.85)",bordercolor='#666',borderwidth=1,borderpad=2,opacity=1,showarrow=False))
    xs=[p[0] for p in pos.values()];ys=[p[1] for p in pos.values()];pad=7
    fig.update_layout(annotations=annotations,showlegend=False,plot_bgcolor='white',margin=dict(t=20,b=20,l=20,r=20),xaxis=dict(visible=False,range=[min(xs)-pad,max(xs)+pad]),yaxis=dict(visible=False,range=[min(ys)-pad,max(ys)+pad]),height=900)
    return fig,G,pos

def renderizar_red_matplotlib(G, pos, titulo="Red", width_in=18, height_in=10):
    fig, ax = plt.subplots(figsize=(width_in, height_in), dpi=150);ax.set_facecolor('white');fig.patch.set_facecolor('white')
    for u,v,data in G.edges(data=True):
        if u in pos and v in pos:w=data.get('weight',1);lw=min(0.5+w*0.05,3.0);ax.plot([pos[u][0],pos[v][0]],[pos[u][1],pos[v][1]],color='#cccccc',linewidth=lw,alpha=0.5,zorder=1)
    for n,d in G.nodes(data=True):
        if n not in pos:continue
        x,y=pos[n];t=d.get('type','palabra');color=d.get('color','#555555');w=float(d.get('weight',1) or 1);label=str(d.get('label',n))
        if t=='candidato':ax.scatter(x,y,c=color,s=800,zorder=5,edgecolors='white',linewidths=3);ax.annotate(label,(x,y),fontsize=16,fontweight='bold',ha='center',va='top',xytext=(0,-22),textcoords='offset points',color='black',zorder=6)
        elif t=='autor':ax.scatter(x,y,c=color,s=80+10*math.log1p(w),zorder=3,edgecolors='white',linewidths=1,marker='D');ax.annotate(label,(x,y),fontsize=13,ha='center',va='center',bbox=dict(boxstyle='round,pad=0.15',facecolor='black',alpha=0.8,edgecolor=color,linewidth=0.5),color='white',zorder=4)
        else:ax.scatter(x,y,c=color,s=60+12*math.log1p(w),zorder=3,edgecolors='white',linewidths=1);ax.annotate(label,(x,y),fontsize=12,ha='center',va='center',bbox=dict(boxstyle='round,pad=0.15',facecolor='black',alpha=0.85,edgecolor=color,linewidth=0.5),color='white',zorder=4)
    ax.axis('off');ax.set_aspect('equal',adjustable='datalim');plt.tight_layout(pad=0.3)
    buf=io.BytesIO();plt.savefig(buf,format='png',dpi=150,bbox_inches='tight',facecolor='white',edgecolor='none');plt.close(fig);buf.seek(0);return buf.read()

def exportar_red_png(fig_plotly, G, pos, titulo="Red", width=2400, height=1200):
    if KALEIDO_OK:
        try:
            img=fig_plotly.to_image(format="png",width=width,height=height,scale=2)
            if img and len(img)>1000:return img
        except:pass
    return renderizar_red_matplotlib(G,pos,titulo,width_in=width/150,height_in=height/150)

def color_func_shades(hex_color):
    base=_hex_to_rgb(hex_color)
    def _f(word,font_size,position,orientation,random_state=None,**kwargs):
        fb=random.uniform(0.10,0.50);r,g,b=_mezclar_con_blanco(base,fb);return f"rgb({r},{g},{b})"
    return _f

def generar_nube_palabras(df_exploded, candidato):
    sub = df_exploded[df_exploded['Candidatos_Detectados']==candidato]
    if sub.empty: return None
    freq_map = obtener_tokens_nube(sub['Contenido de la publicación'], blacklist_extra=list(CANDIDATOS_CONFIG[candidato]['blacklist']))
    if not freq_map: return None
    top_freq = dict(Counter(freq_map).most_common(100))
    wc = WordCloud(width=850,height=480,background_color='white',max_words=100,prefer_horizontal=0.85,relative_scaling=0.45,min_font_size=11,collocations=False,normalize_plurals=False).generate_from_frequencies(top_freq)
    wc.recolor(color_func=color_func_shades(COLORES[candidato]),random_state=7)
    fig, ax = plt.subplots(figsize=(8.5, 4.8));ax.imshow(wc, interpolation='bilinear');ax.axis('off')
    fig.suptitle(candidato, fontsize=20, fontweight='bold', color=COLORES[candidato], y=0.98)
    buf = io.BytesIO();plt.savefig(buf, format='png', bbox_inches='tight', dpi=110, facecolor='white');plt.close(fig);buf.seek(0)
    return base64.b64encode(buf.read()).decode('utf-8')

def calcular_kpis_candidatos(df_exploded):
    kpis={}
    if 'followers' in df_exploded.columns and 'fans' in df_exploded.columns:df_exploded['Alcance_Total']=df_exploded['followers'].fillna(0)+df_exploded['fans'].fillna(0)
    else:df_exploded['Alcance_Total']=0
    for cand in CANDIDATOS_CONFIG:
        sub=df_exploded[df_exploded['Candidatos_Detectados']==cand]
        kpis[cand]={'menciones':int(len(sub)),'alcance':int(sub['Alcance_Total'].sum()),'autores_unicos':int(sub['Autor'].nunique()) if 'Autor' in sub.columns else 0,'color':COLORES[cand]}
    return kpis

# ============================================
# HTML + PPTX (funciones completas, idénticas a Colab)
# Se omiten por brevedad - son las mismas del código anterior
# ============================================

def generar_html_final(df,fig_main,fig_inverse,wordclouds_b64,ejes_tematicos,kpis_ordenados):
    df_exp=df.explode('Candidatos_Detectados').dropna(subset=['Candidatos_Detectados'])
    if 'followers' in df.columns and 'fans' in df.columns:df_exp['Alcance_Total']=df_exp['followers'].fillna(0)+df_exp['fans'].fillna(0)
    else:df_exp['Alcance_Total']=0
    tpc={}
    for cand,eje in ejes_tematicos.items():
        if not eje:continue
        for i in range(1,6):
            tk=f'tema_{i}'
            if tk in eje and eje[tk]:
                n=eje[tk].get('nombre','').strip()
                if n:
                    if n not in tpc:tpc[n]=[]
                    tpc[n].append(cand)
    tc={t for t,c in tpc.items() if len(c)>=2}
    kh='<div class="kpis-row">'
    for cand,data in kpis_ordenados.items():
        color=data['color'];r,g,b=_hex_to_rgb(color)
        kh+=f"""<div class="kpi-card"><div class="kpi-top-bar" style="background:linear-gradient(135deg,{color},rgba({r},{g},{b},0.6))"></div>
        <div class="kpi-body"><div class="kpi-avatar" style="background:linear-gradient(135deg,{color},rgba({r},{g},{b},0.7))">{CANDIDATOS_CONFIG[cand]['nombre_corto'][0]}</div>
        <div class="kpi-name">{CANDIDATOS_CONFIG[cand]['nombre_corto']}</div>
        <div class="kpi-stats-row"><div class="kpi-stat-item"><div class="kpi-stat-num" style="color:{color}">{data['menciones']:,}</div><div class="kpi-stat-label">Menciones</div></div>
        <div class="kpi-stat-divider"></div><div class="kpi-stat-item"><div class="kpi-stat-num" style="color:{color}">{data['alcance']:,}</div><div class="kpi-stat-label">Alcance</div></div>
        <div class="kpi-stat-divider"></div><div class="kpi-stat-item"><div class="kpi-stat-num" style="color:{color}">{data['autores_unicos']:,}</div><div class="kpi-stat-label">Autores</div></div>
        </div></div></div>"""
    kh+='</div>'
    eh='<div class="ejes-row">'
    for cand,eje in ejes_tematicos.items():
        if not eje:continue
        color=COLORES[cand];tp=eje.get('total_posts',0);ti=eje.get('total_interacciones',0);tl=""
        for i in range(1,6):
            tk=f'tema_{i}'
            if tk in eje and eje[tk]:
                tema=eje[tk];nombre=tema.get('nombre','N/A');posts=tema.get('posts',0);inter=tema.get('interacciones',0)
                rk={1:'🥇',2:'🥈',3:'🥉',4:'4️⃣',5:'5️⃣'}.get(i,'•');sb='<span class="shared-badge" title="Compartido">🔗</span>' if nombre in tc else ''
                cls="eje-tema eje-tema-top" if i==1 else "eje-tema";sty=f'color:{color};font-weight:700;' if i==1 else ''
                tl+=f'<div class="{cls}"><span class="eje-rank">{rk}</span><span class="eje-tema-name" style="{sty}">{nombre}</span>{sb}<span class="eje-tema-stats">{posts:,}p · {inter:,}i</span></div>'
        eh+=f"""<div class="eje-col" style="border-top:4px solid {color}"><div class="eje-col-header" style="color:{color}">{CANDIDATOS_CONFIG[cand]['nombre_corto']}</div><div class="eje-col-totals">{tp:,} posts · {ti:,} int.</div>{tl}</div>"""
    eh+='</div>'
    if tc:
        cd="".join([f'<span class="shared-tag">🔗 {t}: {", ".join([CANDIDATOS_CONFIG[c]["nombre_corto"] for c in tpc[t]])}</span>' for t in sorted(tc)])
        eh+=f'<div class="shared-legend"><strong>Temas compartidos:</strong><br><div class="shared-tags-row">{cd}</div></div>'
    wh='<div class="wordclouds-grid">'
    for cand in kpis_ordenados.keys():
        b64=wordclouds_b64.get(cand)
        if b64:wh+=f'<div class="wordcloud-item"><img src="data:image/png;base64,{b64}" style="width:100%;border-radius:10px"></div>'
    wh+='</div>'
    red_col='Grupo de dominio' if 'Grupo de dominio' in df.columns else ('Fuente' if 'Fuente' in df.columns else None)
    th=""
    for cand in kpis_ordenados.keys():
        color=COLORES[cand];sub=df_exp[df_exp['Candidatos_Detectados']==cand]
        if 'Autor' not in sub.columns:continue
        agg_dict={'Alcance_Total':'sum','Contenido de la publicación':'count'}
        if red_col:agg_dict[red_col]=lambda x:x.mode()[0] if not x.mode().empty else 'N/A'
        dfa=sub.groupby('Autor').agg(agg_dict).reset_index().sort_values('Alcance_Total',ascending=False).head(10)
        filas=""
        for _,row in dfa.iterrows():
            autor_nombre=str(row['Autor']);red_valor=str(row[red_col]) if red_col and red_col in row.index else 'N/A'
            alcance_valor=int(row['Alcance_Total']);posts_count=int(row['Contenido de la publicación'])
            alcance_display=f"{alcance_valor:,}" if alcance_valor>0 else f"0 ({posts_count} Post{'s' if posts_count!=1 else ''})"
            filas+=f'<tr><td class="auth-td auth-td-autor" title="@{autor_nombre}">@{autor_nombre}</td><td class="auth-td auth-td-red">{red_valor}</td><td class="auth-td auth-td-alcance" style="color:{color}">{alcance_display}</td></tr>'
        th+=f"""<div class="author-card" style="border-top:4px solid {color}"><h3 class="author-card-title" style="color:{color}">👥 {cand}</h3><table class="author-table"><thead><tr><th class="auth-th" style="width:45%">Autor</th><th class="auth-th" style="width:30%">Red</th><th class="auth-th auth-th-right" style="width:25%">Alcance</th></tr></thead><tbody>{filas}</tbody></table></div>"""
    hashtags=re.findall(r'#\w+'," ".join(df['Contenido de la publicación'].astype(str)).lower());hashtags=[h for h in hashtags if h not in ['#colombiadecide2026','#debate','#colombia']];top_tags=Counter(hashtags).most_common(12)
    html=f"""<!DOCTYPE html><html lang="es"><head><meta charset="UTF-8"><title>Reporte Kreab 2026</title>
<style>body{{font-family:'Segoe UI',system-ui,Arial,sans-serif;background:#f3f5f8;margin:0;padding:18px;color:#101828}}.container{{max-width:1400px;margin:0 auto}}.header{{background:linear-gradient(135deg,#8B0000,#111827);color:white;padding:28px;border-radius:14px;margin-bottom:22px;box-shadow:0 12px 30px rgba(16,24,40,0.22)}}.section-title{{border-left:5px solid #8B0000;padding-left:12px;font-size:22px;margin:40px 0 16px;font-weight:800;color:#111827}}.card-grid{{display:grid;grid-template-columns:repeat(auto-fit,minmax(200px,1fr));gap:12px;margin-bottom:16px}}.card{{background:white;padding:16px;border-radius:12px;text-align:center;box-shadow:0 8px 20px rgba(16,24,40,0.05)}}.number{{font-size:30px;font-weight:900;color:#8B0000;display:block;margin-bottom:3px}}.label{{font-size:11px;color:#667085;text-transform:uppercase;letter-spacing:0.8px;font-weight:700}}.kpis-row{{display:flex;flex-wrap:wrap;gap:14px;margin:16px 0}}.kpi-card{{flex:1;min-width:200px;background:#fff;border-radius:14px;overflow:hidden;box-shadow:0 8px 28px rgba(16,24,40,0.08)}}.kpi-top-bar{{height:6px;width:100%}}.kpi-body{{padding:16px 14px 14px;text-align:center}}.kpi-avatar{{width:42px;height:42px;border-radius:50%;color:white;font-size:20px;font-weight:900;display:flex;align-items:center;justify-content:center;margin:0 auto 8px;box-shadow:0 4px 12px rgba(0,0,0,0.15)}}.kpi-name{{font-size:13px;font-weight:800;color:#111827;margin-bottom:12px;text-transform:uppercase;letter-spacing:0.4px}}.kpi-stats-row{{display:flex;justify-content:center;align-items:center}}.kpi-stat-item{{flex:1;text-align:center;padding:0 6px}}.kpi-stat-num{{display:block;font-size:17px;font-weight:900;line-height:1.2}}.kpi-stat-label{{display:block;font-size:9px;color:#667085;text-transform:uppercase;letter-spacing:0.5px;font-weight:700;margin-top:2px}}.kpi-stat-divider{{width:1px;height:30px;background:#e5e7eb}}.ejes-row{{display:flex;flex-wrap:nowrap;gap:12px;margin:16px 0;overflow-x:auto;padding-bottom:8px}}.eje-col{{flex:1;min-width:210px;background:white;border-radius:12px;padding:12px;box-shadow:0 8px 22px rgba(16,24,40,0.07)}}.eje-col-header{{font-size:12px;font-weight:800;text-transform:uppercase;letter-spacing:0.3px;margin-bottom:3px}}.eje-col-totals{{font-size:9px;color:#888;margin-bottom:10px;padding-bottom:6px;border-bottom:1px solid #f0f0f0}}.eje-tema{{display:flex;align-items:center;gap:5px;padding:4px 5px;border-radius:5px;margin-bottom:2px;background:#fafafa;border:1px solid #f5f5f5}}.eje-tema-top{{background:#f0f7ff;border:1px solid #dbeafe}}.eje-rank{{font-size:11px;flex-shrink:0}}.eje-tema-name{{font-size:10.5px;font-weight:600;color:#333;flex:1;line-height:1.25}}.eje-tema-stats{{font-size:8px;color:#999;white-space:nowrap;flex-shrink:0}}.shared-badge{{font-size:10px;cursor:help;flex-shrink:0}}.shared-legend{{background:#fffbeb;border:1px solid #fde68a;border-radius:10px;padding:12px 16px;margin-top:12px;font-size:12px;color:#92400e}}.shared-tags-row{{display:flex;flex-wrap:wrap;gap:8px;margin-top:6px}}.shared-tag{{background:#fef3c7;color:#92400e;padding:4px 10px;border-radius:6px;font-size:11px;font-weight:600}}.wordclouds-grid{{display:grid;grid-template-columns:repeat(auto-fit,minmax(340px,1fr));gap:16px;margin:16px 0}}.wordcloud-item{{background:white;padding:12px;border-radius:12px;box-shadow:0 10px 28px rgba(16,24,40,0.07)}}.network-box{{background:white;padding:8px;border-radius:14px;box-shadow:0 12px 30px rgba(16,24,40,0.07);margin-bottom:22px}}.authors-grid{{display:flex;flex-wrap:wrap;gap:16px;margin:16px 0}}.author-card{{flex:1;min-width:250px;max-width:100%;background:white;border-radius:10px;padding:14px;box-shadow:0 8px 20px rgba(16,24,40,0.07)}}.author-card-title{{margin-top:0;margin-bottom:10px;font-size:15px;font-weight:700}}.author-table{{width:100%;border-collapse:collapse;font-size:12px;table-layout:fixed}}.auth-th{{padding:8px 6px;text-align:left;background:#f9fafb;font-weight:700;color:#344054;font-size:11px;border-bottom:2px solid #e5e7eb}}.auth-th-right{{text-align:right}}.auth-td{{padding:8px 6px;border-bottom:1px solid #f0f0f0;vertical-align:middle}}.auth-td-autor{{overflow:hidden;text-overflow:ellipsis;white-space:nowrap;font-weight:500;color:#111827}}.auth-td-red{{overflow:hidden;text-overflow:ellipsis;white-space:nowrap;font-size:11px;color:#667085}}.auth-td-alcance{{text-align:right;font-weight:700;white-space:nowrap}}.tag-cloud{{display:flex;flex-wrap:wrap;gap:8px}}.tag{{background:#eef2f6;color:#344054;padding:6px 10px;border-radius:999px;font-size:12px;font-weight:700}}</style></head><body>
<div class="container"><div class="header"><h1 style="margin:0">🗳️ Inteligencia Electoral #ColombiaDecide2026</h1><p style="margin:5px 0 0;opacity:0.9">KREAB COLOMBIA</p></div>
<div class="card-grid"><div class="card"><span class="number">{len(df):,}</span><span class="label">Publicaciones</span></div><div class="card"><span class="number">{len(df_exp):,}</span><span class="label">Menciones</span></div><div class="card"><span class="number">{int(df_exp['Alcance_Total'].sum()):,}</span><span class="label">Alcance Total</span></div></div>
<h2 class="section-title">📊 KPIs por Candidato</h2>{kh}
<h2 class="section-title">📌 Ejes Temáticos por Candidato</h2><p style="color:#667085;margin-bottom:10px;font-size:12px">🔗 = tema compartido.</p>{eh}
<h2 class="section-title">#️⃣ Hashtags Principales</h2><div class="tag-cloud">{''.join([f'<span class="tag">{t} ({c})</span>' for t,c in top_tags])}</div>
<h2 class="section-title" style="font-size:28px">🕸️ Red de Narrativas</h2><div class="network-box">{fig_main.to_html(full_html=False,include_plotlyjs='cdn')}</div>
<h2 class="section-title" style="font-size:28px">🔄 Narrativas Compartidas</h2><div class="network-box">{fig_inverse.to_html(full_html=False,include_plotlyjs=False)}</div>
<h2 class="section-title">☁️ Nubes de Palabras</h2>{wh}
<h2 class="section-title">👥 Autores Influyentes</h2><div class="authors-grid">{th}</div></div></body></html>"""
    return html

def _hex_to_rgbcolor(hc):hc=hc.lstrip('#');return RGBColor(int(hc[0:2],16),int(hc[2:4],16),int(hc[4:6],16))
def _add_rr(sl,l,t,w,h,f='#FFFFFF',b=None,bw=1.0,cs=None):
    s=sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h);s.fill.solid();s.fill.fore_color.rgb=_hex_to_rgbcolor(f)
    if b:s.line.color.rgb=_hex_to_rgbcolor(b);s.line.width=Pt(bw)
    else:s.line.fill.background()
    s.adjustments[0]=cs if cs is not None else 0.05;return s
def _add_tb(sl,l,t,w,h,tx,fs=12,fc='#111827',bd=False,al=PP_ALIGN.LEFT,fn='Segoe UI'):
    tb=sl.shapes.add_textbox(l,t,w,h);tf=tb.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.text=tx;p.font.size=Pt(fs);p.font.color.rgb=_hex_to_rgbcolor(fc);p.font.bold=bd;p.font.name=fn;p.alignment=al;return tb
def _set_bg(sl,c='#F3F5F8'):f=sl.background.fill;f.solid();f.fore_color.rgb=_hex_to_rgbcolor(c)

def generar_powerpoint(df,fig_main,fig_inverse,wordclouds_b64,ejes_tematicos,kpis_ordenados,png1,png2):
    if not PPTX_AVAILABLE:return None
    prs=Presentation();prs.slide_width=Inches(13.333);prs.slide_height=Inches(7.5);SW=Inches(13.333);M=Inches(0.4);CW=SW-2*M;CD='#FFFFFF';DK='#111827';MT='#667085';AC='#8B0000'
    df_exp=df.explode('Candidatos_Detectados').dropna(subset=['Candidatos_Detectados'])
    if 'followers' in df.columns and 'fans' in df.columns:df_exp['Alcance_Total']=df_exp['followers'].fillna(0)+df_exp['fans'].fillna(0)
    else:df_exp['Alcance_Total']=0
    rc='Grupo de dominio' if 'Grupo de dominio' in df.columns else ('Fuente' if 'Fuente' in df.columns else None)
    ht=re.findall(r'#\w+'," ".join(df['Contenido de la publicación'].astype(str)).lower());ht=[h for h in ht if h not in ['#colombiadecide2026','#debate','#colombia']];tt=Counter(ht).most_common(12)
    tpc2={}
    for cn,ej in ejes_tematicos.items():
        if not ej:continue
        for i in range(1,6):
            tk=f'tema_{i}'
            if tk in ej and ej[tk]:
                nm=ej[tk].get('nombre','').strip()
                if nm:
                    if nm not in tpc2:tpc2[nm]=[]
                    tpc2[nm].append(cn)
    tcs={t for t,c in tpc2.items() if len(c)>=2};nc=len(kpis_ordenados)
    s1=prs.slides.add_slide(prs.slide_layouts[6]);_set_bg(s1);hd=s1.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0),Inches(0),SW,Inches(0.9));hd.fill.solid();hd.fill.fore_color.rgb=_hex_to_rgbcolor(AC);hd.line.fill.background()
    _add_tb(s1,Inches(0.5),Inches(0.12),Inches(10),Inches(0.35),"🗳️ Inteligencia Electoral #ColombiaDecide2026",22,'#FFFFFF',True);_add_tb(s1,Inches(0.5),Inches(0.50),Inches(10),Inches(0.3),"KREAB COLOMBIA",11,'#FFD0D0',False)
    tp2=len(df);tm=len(df_exp);ta=int(df_exp['Alcance_Total'].sum());gy=Inches(1.05);gw=Inches(3.8);gg=Inches(0.25)
    for idx,(v,lb) in enumerate([(tp2,"Publicaciones"),(tm,"Menciones"),(ta,"Alcance Total")]):x=M+idx*(gw+gg);_add_rr(s1,x,gy,gw,Inches(0.55),CD,'#E5E7EB',0.5);_add_tb(s1,x+Inches(0.15),gy+Inches(0.05),Inches(2),Inches(0.28),f"{v:,}",18,AC,True);_add_tb(s1,x+Inches(2.2),gy+Inches(0.12),Inches(1.5),Inches(0.25),lb,9,MT,True)
    _add_tb(s1,M,Inches(1.72),Inches(8),Inches(0.3),"📊 KPIs por Candidato",14,DK,True);ky=Inches(2.02);kw=(CW-Inches(0.15)*(nc-1))/nc
    for idx,(cn,dt) in enumerate(kpis_ordenados.items()):x=M+idx*(kw+Inches(0.15));co=dt['color'];_add_rr(s1,x,ky,kw,Inches(0.90),CD,co,1.5);br=s1.shapes.add_shape(MSO_SHAPE.RECTANGLE,x,ky,kw,Inches(0.06));br.fill.solid();br.fill.fore_color.rgb=_hex_to_rgbcolor(co);br.line.fill.background();_add_tb(s1,x+Inches(0.08),ky+Inches(0.10),kw-Inches(0.16),Inches(0.22),CANDIDATOS_CONFIG[cn]['nombre_corto'],10,co,True,PP_ALIGN.CENTER);st=f"Menciones: {dt['menciones']:,}  |  Alcance: {dt['alcance']:,}  |  Autores: {dt['autores_unicos']:,}";_add_tb(s1,x+Inches(0.05),ky+Inches(0.38),kw-Inches(0.1),Inches(0.45),st,8,'#555555',False,PP_ALIGN.CENTER)
    _add_tb(s1,M,Inches(3.05),Inches(12),Inches(0.3),"📌 Ejes Temáticos por Candidato",12,DK,True);ey=Inches(3.35);ew=(CW-Inches(0.12)*(nc-1))/nc
    for idx,(cn,ej) in enumerate(ejes_tematicos.items()):
        if not ej:continue
        x=M+idx*(ew+Inches(0.12));co=COLORES[cn];_add_rr(s1,x,ey,ew,Inches(3.10),CD,co,1.0);br=s1.shapes.add_shape(MSO_SHAPE.RECTANGLE,x,ey,ew,Inches(0.05));br.fill.solid();br.fill.fore_color.rgb=_hex_to_rgbcolor(co);br.line.fill.background()
        _add_tb(s1,x+Inches(0.06),ey+Inches(0.07),ew-Inches(0.12),Inches(0.20),CANDIDATOS_CONFIG[cn]['nombre_corto'],9,co,True,PP_ALIGN.CENTER);_add_tb(s1,x+Inches(0.06),ey+Inches(0.27),ew-Inches(0.12),Inches(0.14),f"{ej.get('total_posts',0):,} posts · {ej.get('total_interacciones',0):,} int.",6,MT,False,PP_ALIGN.CENTER)
        tsy=ey+Inches(0.46)
        for i in range(1,6):
            tk=f'tema_{i}'
            if tk in ej and ej[tk]:tm2=ej[tk];nm=tm2.get('nombre','N/A');ps=tm2.get('posts',0);it=tm2.get('interacciones',0);rk={1:'#1',2:'#2',3:'#3',4:'#4',5:'#5'}.get(i,'•');sm=" 🔗" if nm in tcs else "";ty=tsy+(i-1)*Inches(0.48);tbg='#F0F7FF' if i==1 else '#FAFAFA';_add_rr(s1,x+Inches(0.05),ty,ew-Inches(0.10),Inches(0.42),tbg,'#E8E8E8',0.5,0.08);_add_tb(s1,x+Inches(0.08),ty+Inches(0.02),ew-Inches(0.20),Inches(0.20),f"{rk} {nm}{sm}",8 if i==1 else 7,co if i==1 else '#333333',i==1);_add_tb(s1,x+Inches(0.08),ty+Inches(0.22),ew-Inches(0.20),Inches(0.16),f"{ps:,} posts · {it:,} int.",6,'#999999',False)
    _add_tb(s1,M,Inches(6.55),Inches(4),Inches(0.3),"#️⃣ Hashtags Principales",14,DK,True);tgt="   ".join([f"{t} ({c})" for t,c in tt]);_add_rr(s1,M,Inches(6.85),CW,Inches(0.45),CD,'#E5E7EB',0.5);_add_tb(s1,M+Inches(0.15),Inches(6.90),CW-Inches(0.3),Inches(0.35),tgt,9,'#344054',True)
    rx=M;ry=Inches(0.85);rw=CW;rh=Inches(6.45)
    s2=prs.slides.add_slide(prs.slide_layouts[6]);_set_bg(s2);h2=s2.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0),Inches(0),SW,Inches(0.7));h2.fill.solid();h2.fill.fore_color.rgb=_hex_to_rgbcolor(AC);h2.line.fill.background();_add_tb(s2,Inches(0.5),Inches(0.15),Inches(10),Inches(0.35),"🕸️ Red de Narrativas (ForceAtlas)",22,'#FFFFFF',True);_add_rr(s2,rx,ry,rw,rh,CD,'#D1D5DB',1.0);s2.shapes.add_picture(io.BytesIO(png1),rx+Inches(0.1),ry+Inches(0.1),rw-Inches(0.2),rh-Inches(0.2))
    s3=prs.slides.add_slide(prs.slide_layouts[6]);_set_bg(s3);h3=s3.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0),Inches(0),SW,Inches(0.7));h3.fill.solid();h3.fill.fore_color.rgb=_hex_to_rgbcolor(AC);h3.line.fill.background();_add_tb(s3,Inches(0.5),Inches(0.15),Inches(10),Inches(0.35),"🔄 Narrativas Compartidas (Red Inversa)",22,'#FFFFFF',True);_add_rr(s3,rx,ry,rw,rh,CD,'#D1D5DB',1.0);s3.shapes.add_picture(io.BytesIO(png2),rx+Inches(0.1),ry+Inches(0.1),rw-Inches(0.2),rh-Inches(0.2))
    s4=prs.slides.add_slide(prs.slide_layouts[6]);_set_bg(s4);h4=s4.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0),Inches(0),SW,Inches(0.7));h4.fill.solid();h4.fill.fore_color.rgb=_hex_to_rgbcolor(AC);h4.line.fill.background();_add_tb(s4,Inches(0.5),Inches(0.15),Inches(10),Inches(0.35),"☁️ Nubes de Palabras",20,'#FFFFFF',True)
    wi=[(c,wordclouds_b64[c]) for c in kpis_ordenados.keys() if wordclouds_b64.get(c)];nw=len(wi)
    if nw>0:
        if nw<=3:cols,rows=nw,1
        elif nw<=4:cols,rows=2,2
        else:cols,rows=3,2
        way=Inches(0.85);wah=Inches(6.3);cg=Inches(0.15);cw2=(CW-cg*(cols-1))/cols;ch2=(wah-cg*(rows-1))/rows
        for i,(cn,b64) in enumerate(wi):r=i//cols;ci=i%cols;cx=M+ci*(cw2+cg);cy=way+r*(ch2+cg);_add_rr(s4,cx,cy,cw2,ch2,CD,COLORES.get(cn,'#999'),1.5);pd2=Inches(0.08);s4.shapes.add_picture(io.BytesIO(base64.b64decode(b64)),cx+pd2,cy+pd2,cw2-2*pd2,ch2-2*pd2)
    s5=prs.slides.add_slide(prs.slide_layouts[6]);_set_bg(s5);h5=s5.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0),Inches(0),SW,Inches(0.7));h5.fill.solid();h5.fill.fore_color.rgb=_hex_to_rgbcolor(AC);h5.line.fill.background();_add_tb(s5,Inches(0.5),Inches(0.15),Inches(10),Inches(0.35),"👥 Autores Influyentes y Alcance",20,'#FFFFFF',True)
    if nc<=3:tc2,tr2=nc,1
    else:tc2,tr2=3,2
    tay=Inches(0.85);tah=Inches(6.3);tg2=Inches(0.15);tcw=(CW-tg2*(tc2-1))/tc2;tch=(tah-tg2*(tr2-1))/tr2
    for idx,cn in enumerate(kpis_ordenados.keys()):
        r_idx=idx//tc2;ci=idx%tc2;cx=M+ci*(tcw+tg2);cy=tay+r_idx*(tch+tg2);co=COLORES[cn];_add_rr(s5,cx,cy,tcw,tch,CD,co,1.5);br=s5.shapes.add_shape(MSO_SHAPE.RECTANGLE,cx,cy,tcw,Inches(0.05));br.fill.solid();br.fill.fore_color.rgb=_hex_to_rgbcolor(co);br.line.fill.background();_add_tb(s5,cx+Inches(0.1),cy+Inches(0.08),tcw-Inches(0.2),Inches(0.25),f"👥 {CANDIDATOS_CONFIG[cn]['nombre_corto']}",11,co,True)
        sub=df_exp[df_exp['Candidatos_Detectados']==cn]
        if 'Autor' not in sub.columns or sub.empty:_add_tb(s5,cx+Inches(0.1),cy+Inches(0.4),tcw-Inches(0.2),Inches(0.2),"Sin datos",9,MT,False);continue
        agg_dict={'Alcance_Total':'sum','Contenido de la publicación':'count'}
        if rc:agg_dict[rc]=lambda x:x.mode()[0] if not x.mode().empty else 'N/A'
        dfa=sub.groupby('Autor').agg(agg_dict).reset_index().sort_values('Alcance_Total',ascending=False).head(10)
        col_pad=Inches(0.08);usable_w=tcw-2*col_pad;col_autor_w=usable_w*0.62;col_red_w=usable_w*0.18;col_alc_w=usable_w*0.20;col_autor_x=cx+col_pad;col_alc_x=cx+tcw-col_pad-col_alc_w;col_red_x=col_alc_x-col_red_w
        header_y=cy+Inches(0.34);_add_tb(s5,col_autor_x,header_y,col_autor_w,Inches(0.14),"Autor",6.5,MT,True);_add_tb(s5,col_red_x,header_y,col_red_w,Inches(0.14),"Red",6.5,MT,True);_add_tb(s5,col_alc_x,header_y,col_alc_w,Inches(0.14),"Alcance",6.5,MT,True,PP_ALIGN.RIGHT)
        for ri,(_,row) in enumerate(dfa.iterrows()):
            ry2=cy+Inches(0.50)+ri*Inches(0.22);_add_tb(s5,col_autor_x,ry2,col_autor_w,Inches(0.20),f"@{str(row['Autor'])}",6,'#333333',False)
            red_valor=str(row[rc]) if rc and rc in row.index else 'N/A';red_abrev=red_valor if len(red_valor)<=10 else red_valor[:9]+"…";_add_tb(s5,col_red_x,ry2,col_red_w,Inches(0.20),red_abrev,5.5,MT,False)
            alcance_val=int(row['Alcance_Total']);posts_val=int(row['Contenido de la publicación']);alc_text=f"{alcance_val:,}" if alcance_val>0 else f"0 ({posts_val} Post{'s' if posts_val!=1 else ''})";_add_tb(s5,col_alc_x,ry2,col_alc_w,Inches(0.20),alc_text,6,co,True,PP_ALIGN.RIGHT)
    buf=io.BytesIO();prs.save(buf);buf.seek(0);return buf.getvalue()

# ============================================
# PROCESAMIENTO COMPLETO (con cache en session_state)
# ============================================

def procesar_datos(uploaded_file):
    """Ejecuta todo el análisis y guarda en session_state."""
    df = pd.read_excel(uploaded_file);df.columns = df.columns.str.strip()
    for col_name in ['Grupo de dominio', 'Fuente']:
        if col_name in df.columns:df[col_name]=df[col_name].astype(str).str.replace('Twitter','X',case=False,regex=False)

    df['Candidatos_Detectados']=df['Contenido de la publicación'].apply(identificar_candidatos)
    df_exploded=df.explode('Candidatos_Detectados').dropna(subset=['Candidatos_Detectados'])

    client=obtener_cliente_openai()
    ejes_tematicos={}
    progress=st.progress(0,text="Analizando...")
    cands_list=list(CANDIDATOS_CONFIG.keys())
    for i,cand in enumerate(cands_list):
        progress.progress((i+1)/(len(cands_list)+4),text=f"📌 {CANDIDATOS_CONFIG[cand]['nombre_corto']}...")
        if client:eje=analizar_ejes_tematicos_con_ia(df_exploded,cand,client)
        else:
            sub=df_exploded[df_exploded['Candidatos_Detectados']==cand];bl=CANDIDATOS_CONFIG[cand]['blacklist'];toks=obtener_tokens_filtrados(sub['Contenido de la publicación'],bl,incluir_hashtags=True);freq=Counter(toks).most_common(30);ci=None
            for col in ['Interacciones totales','interacciones totales','Interacciones','interactions']:
                if col in sub.columns:ci=col;break
            ti=int(sub[ci].fillna(0).sum()) if ci else len(sub);eje=generar_analisis_fallback(sub,freq,cand,len(sub),ti,ci)
        ejes_tematicos[cand]=eje

    progress.progress(0.7,text="🔄 Normalizando...")
    ejes_tematicos=normalizar_temas_entre_candidatos(ejes_tematicos,client)
    if 'followers' in df_exploded.columns and 'fans' in df_exploded.columns:df_exploded['Alcance_Total']=df_exploded['followers'].fillna(0)+df_exploded['fans'].fillna(0)
    else:df_exploded['Alcance_Total']=0
    kpis=calcular_kpis_candidatos(df_exploded);kpis_ordenados=ordenar_kpis_por_menciones(kpis);ejes_tematicos=ordenar_ejes_por_menciones_interacciones(ejes_tematicos,kpis)

    progress.progress(0.8,text="⚙️ Redes...")
    fig1,G1,pos1=crear_red_principal(df_exploded);fig2,G2,pos2=crear_red_inversa_gephi(df_exploded)

    progress.progress(0.85,text="📸 PNGs...")
    png1=exportar_red_png(fig1,G1,pos1,"Red",2400,1200);png2=exportar_red_png(fig2,G2,pos2,"Inv",2400,1200)

    progress.progress(0.9,text="☁️ Nubes...")
    wordclouds={}
    for cand in CANDIDATOS_CONFIG:wordclouds[cand]=generar_nube_palabras(df_exploded,cand)

    progress.progress(0.95,text="📝 Reportes...")
    html=generar_html_final(df,fig1,fig2,wordclouds,ejes_tematicos,kpis_ordenados)
    pptx_bytes=generar_powerpoint(df,fig1,fig2,wordclouds,ejes_tematicos,kpis_ordenados,png1,png2)
    progress.progress(1.0,text="✅ ¡Listo!")

    # Guardar todo en session_state
    st.session_state.results = {
        'html': html, 'pptx_bytes': pptx_bytes, 'fig1': fig1, 'fig2': fig2,
        'kpis_ordenados': kpis_ordenados, 'wordclouds': wordclouds,
        'file_name': uploaded_file.name
    }

# ============================================
# APP PRINCIPAL
# ============================================

def run_app():
    with st.sidebar:
        st.markdown("### 🤖📊 PPTX JC Intelligence")
        st.markdown("---")
        uploaded_file = st.file_uploader("📂 Sube tu Excel", type=['xlsx','xls'])
        st.markdown("---")
        st.markdown(f"- OpenAI: {'✅' if OPENAI_AVAILABLE else '❌'}")
        st.markdown(f"- Kaleido: {'✅' if KALEIDO_OK else '⚠️'}")
        st.markdown(f"- PPTX: {'✅' if PPTX_AVAILABLE else '❌'}")
        st.markdown("---")
        if st.button("🚪 Cerrar sesión"):
            st.session_state.authenticated=False;st.session_state.pop('results',None);st.rerun()

    st.markdown("""
    <div style="background:linear-gradient(135deg,#8B0000,#111827);color:white;padding:28px;border-radius:14px;margin-bottom:22px">
        <h1 style="margin:0">🤖📊 Inteligencia Electoral</h1>
        <p style="margin:5px 0 0;opacity:0.9">Análisis de Redes y Narrativas_pptx generator_jc 🤖😺</p>
    </div>
    """, unsafe_allow_html=True)

    if not uploaded_file:
        st.session_state.pop('results', None)
        st.info("👈 Sube un archivo Excel (.xlsx) en la barra lateral para comenzar.")
        return

    # Solo procesar si es un archivo nuevo o no hay resultados
    needs_processing = False
    if 'results' not in st.session_state:
        needs_processing = True
    elif st.session_state.results.get('file_name') != uploaded_file.name:
        needs_processing = True

    if needs_processing:
        procesar_datos(uploaded_file)

    # Mostrar resultados desde cache
    res = st.session_state.results
    html = res['html']
    pptx_bytes = res['pptx_bytes']
    fig1 = res['fig1']
    fig2 = res['fig2']
    kpis_ordenados = res['kpis_ordenados']
    wordclouds = res['wordclouds']

    st.markdown("---")

    # Descargas (ya no re-ejecutan el análisis)
    col_d1, col_d2, _ = st.columns([1,1,2])
    with col_d1:
        st.download_button("📥 Descargar HTML", html.encode('utf-8'), "Kreab_Reporte_2026.html", "text/html", use_container_width=True)
    with col_d2:
        if pptx_bytes:
            st.download_button("📥 Descargar PPTX", pptx_bytes, "Kreab_Reporte_2026.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation", use_container_width=True)

    # KPIs
    st.markdown("### 📊 KPIs por Candidato")
    cols = st.columns(len(kpis_ordenados))
    for i,(cand,data) in enumerate(kpis_ordenados.items()):
        with cols[i]:
            color=data['color']
            st.markdown(f"""<div style="background:white;border-radius:12px;padding:16px;text-align:center;border-top:4px solid {color};box-shadow:0 4px 12px rgba(0,0,0,0.08)">
                <div style="font-size:13px;font-weight:800;color:{color}">{CANDIDATOS_CONFIG[cand]['nombre_corto']}</div>
                <div style="font-size:22px;font-weight:900;color:{color}">{data['menciones']:,}</div>
                <div style="font-size:9px;color:#667085;text-transform:uppercase">Menciones</div>
                <div style="font-size:14px;font-weight:700;color:{color};margin-top:4px">{data['alcance']:,}</div>
                <div style="font-size:9px;color:#667085;text-transform:uppercase">Alcance</div>
                <div style="font-size:12px;color:#555;margin-top:4px">{data['autores_unicos']:,} autores</div>
            </div>""", unsafe_allow_html=True)

    st.markdown("### 🕸️ Red de Narrativas")
    st.plotly_chart(fig1, use_container_width=True)
    st.markdown("### 🔄 Narrativas Compartidas")
    st.plotly_chart(fig2, use_container_width=True)

    st.markdown("### ☁️ Nubes de Palabras")
    wc_items=[(c,wordclouds[c]) for c in kpis_ordenados.keys() if wordclouds.get(c)]
    if wc_items:
        wc_cols=st.columns(min(3,len(wc_items)))
        for i,(cand,b64) in enumerate(wc_items):
            with wc_cols[i%len(wc_cols)]:
                st.image(f"data:image/png;base64,{b64}", use_container_width=True)

    st.markdown("### 📄 Reporte Completo")
    with st.expander("Ver HTML completo", expanded=False):
        st.components.v1.html(html, height=3000, scrolling=True)

# ============================================
# ENTRY POINT
# ============================================

if __name__ == "__main__":
    if check_password():
        run_app()
